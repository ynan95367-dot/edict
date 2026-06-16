#!/usr/bin/env python3
"""
Full PPT analysis - v2 robust version
Analyzes /Users/bingsen/Desktop/字节分享final-0612.pptx.pptx
Outputs structured JSON for report generation.
"""

import json, sys, os, zipfile
from collections import Counter, defaultdict
from lxml import etree

PPT_PATH = "/Users/bingsen/Desktop/字节分享final-0612.pptx.pptx"

# Shape type name mapping
MSO_TYPE_NAMES = {
    1: "AUTO_SHAPE", 2: "CALLOUT", 3: "CANVAS", 4: "CHART",
    5: "COMMENT", 6: "CONNECTOR", 7: "DIAGRAM", 8: "FREE_FORM",
    9: "GROUP", 10: "INK", 11: "LINE", 12: "LINKED_OLE_OBJECT",
    13: "LINKED_PICTURE", 14: "MEDIA", 15: "OLE_OBJECT", 16: "PICTURE",
    17: "PLACEHOLDER", 18: "RADIO_BUTTON", 19: "SCRIPT_ANCHOR",
    20: "SHAPE_TYPE_MIX", 21: "TABLE", 22: "TEXT_BOX",
    23: "TEXT_EFFECT", 24: "WEB_VIDEO",
}

def get_font_name(run):
    """Safely extract font name from a run."""
    try:
        fn = run.font.name
        return fn if fn else None
    except:
        return None

def is_chinese_font(name):
    if not name:
        return False
    n = name.lower()
    cjk_kw = ["hei","song","fang","kai","ming","yuan","wenquan",
              "pingfang","microsoft yahei","simsun","simhei",
              "noto sans sc","noto serif sc","fangzheng","fz",
              "华文","微软雅黑","黑体","宋体","楷体","仿宋","思源","汉仪","方正"]
    for kw in cjk_kw:
        if kw in n:
            return True
    return False

def main():
    # ZIP check
    if not zipfile.is_zipfile(PPT_PATH):
        print(json.dumps({"error": "Not a valid ZIP/PPTX"}, ensure_ascii=False))
        sys.exit(1)

    from pptx import Presentation
    prs = Presentation(PPT_PATH)

    results = {
        "file_path": PPT_PATH,
        "file_size_bytes": os.path.getsize(PPT_PATH),
        "zip_valid": True,
        "total_slides": len(prs.slides),
        "slides": [],
        "notes": {},
        "fonts_all": [],
        "fonts_chinese": set(),
        "fonts_nonchinese": set(),
        "shape_stats": Counter(),
        "layout_stats": Counter(),
        "image_count": 0,
        "chart_count": 0,
        "table_count": 0,
        "group_count": 0,
        "errors": [],
    }

    for idx, slide in enumerate(prs.slides, 1):
        sd = {
            "num": idx,
            "layout": None,
            "shapes_count": 0,
            "text_blocks": [],
            "text_preview": "",
            "shape_types": [],
            "image_count": 0,
        }
        try:
            # Layout
            try:
                sd["layout"] = slide.slide_layout.name if slide.slide_layout else "Unknown"
            except:
                sd["layout"] = "Unknown"

            # Shapes
            for shape in slide.shapes:
                st_val = shape.shape_type
                if isinstance(st_val, int):
                    st_name = MSO_TYPE_NAMES.get(st_val, f"OTHER({st_val})")
                else:
                    st_name = MSO_TYPE_NAMES.get(st_val.value, f"OTHER({st_val.value})")
                sd["shape_types"].append(st_name)
                results["shape_stats"][st_name] += 1

                # Count specific types
                if st_val == 16:  # PICTURE
                    results["image_count"] += 1
                    sd["image_count"] += 1
                elif st_val == 4:  # CHART
                    results["chart_count"] += 1
                elif st_val == 21:  # TABLE
                    results["table_count"] += 1
                elif st_val == 9:  # GROUP
                    results["group_count"] += 1

                # Text extraction
                try:
                    if shape.has_text_frame:
                        full = shape.text_frame.text.strip()
                        if full:
                            sd["text_blocks"].append(full[:120])
                except:
                    pass

                # Font audit
                try:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                fn = get_font_name(run)
                                if fn:
                                    results["fonts_all"].append(fn)
                                    if is_chinese_font(fn):
                                        results["fonts_chinese"].add(fn)
                                    else:
                                        results["fonts_nonchinese"].add(fn)
                except:
                    pass

            sd["shapes_count"] = len(sd["shape_types"])
            text = "\n".join(sd["text_blocks"])
            sd["text_preview"] = text[:200].replace("\n", " | ") if text else "(empty)"

            # Speaker notes
            try:
                if slide.has_notes_slide:
                    ns = slide.notes_slide
                    if ns and ns.notes_text_frame:
                        nt = ns.notes_text_frame.text.strip()
                        if nt:
                            results["notes"][idx] = nt[:300]
            except:
                pass

            # Layout stats
            results["layout_stats"][sd["layout"]] += 1

        except Exception as e:
            results["errors"].append(f"Slide {idx}: {type(e).__name__}: {e}")

        results["slides"].append(sd)

    # Convert sets to sorted lists
    results["fonts_chinese"] = sorted(results["fonts_chinese"])
    results["fonts_nonchinese"] = sorted(results["fonts_nonchinese"])
    results["font_total_unique"] = len(set(results["fonts_all"]))

    # Convert Counters to dict
    results["shape_stats"] = dict(results["shape_stats"])
    results["layout_stats"] = dict(results["layout_stats"])

    print(json.dumps(results, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    main()
