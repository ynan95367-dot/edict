#!/usr/bin/env python3
"""
Debug script to understand the PPTX structure.
"""
import sys

PPT_PATH = "/Users/bingsen/Desktop/字节分享final-0612.pptx.pptx"

# First, check the XML structure inside the PPTX
import zipfile
from lxml import etree

z = zipfile.ZipFile(PPT_PATH)
print("=== ZIP contents ===")
for name in z.namelist():
    print(f"  {name} ({z.getinfo(name).file_size} bytes)")

# Check presentation.xml for slide list and layout info
if 'ppt/presentation.xml' in z.namelist():
    pres_xml = z.read('ppt/presentation.xml')
    root = etree.fromstring(pres_xml)
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
          'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
    print("\n=== Slide IDs ===")
    for sldId in root.findall('.//p:sldId', ns):
        print(f"  id={sldId.get('id')}, r:id={sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')}")

    # Layout info
    print("\n=== Slide Layout IDs ===")
    for layout in root.findall('.//p:sldLayoutId', ns):
        print(f"  id={layout.get('id')}, r:id={layout.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')}")

# Check actual slide XMLs
print("\n=== FIRST SLIDE XML (sample) ===")
for name in sorted(z.namelist()):
    if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
        xml_content = z.read(name)
        root = etree.fromstring(xml_content)
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
              'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
              'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
        
        # Count shapes
        shapes = root.findall('.//p:sp', ns) + root.findall('.//p:pic', ns) + root.findall('.//p:graphicFrame', ns) + root.findall('.//p:tbl', ns) + root.findall('.//p:grpSp', ns)
        print(f"\n{name}: {len(shapes)} shapes found")
        
        # Show first 500 chars of XML
        text = etree.tostring(root, pretty_print=True).decode()
        text_lines = text.split('\n')
        print('\n'.join(text_lines[:30]))
        
        # Only show first slide
        break

# Also check what the python-pptx shapes look like
print("\n=== python-pptx slide.shapes debug ===")
from pptx import Presentation
prs = Presentation(PPT_PATH)
for idx, slide in enumerate(prs.slides, 1):
    if idx > 3:
        break
    print(f"\nSlide {idx}:")
    try:
        shapes_list = list(slide.shapes)
        print(f"  Number of shapes: {len(shapes_list)}")
        for shape in shapes_list:
            print(f"  Shape: name='{shape.name}', type={shape.shape_type}, type_name={type(shape.shape_type)}")
    except Exception as e:
        print(f"  Error iterating shapes: {type(e).__name__}: {e}")
        import traceback
        traceback.print_exc()

z.close()
