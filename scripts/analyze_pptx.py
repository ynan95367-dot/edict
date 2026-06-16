#!/usr/bin/env python3
"""
字节分享PPT全面审查脚本
分析 /Users/bingsen/Desktop/字节分享final-0612.pptx.pptx
输出结构化审查数据供报告生成
"""

import json
import sys
from collections import Counter, defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPT_PATH = "/Users/bingsen/Desktop/字节分享final-0612.pptx.pptx"

def shape_type_name(st):
    """Map shape type enum to readable name."""
    mapping = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "AutoShape",
        MSO_SHAPE_TYPE.CALLOUT: "Callout",
        MSO_SHAPE_TYPE.CANVAS: "Canvas",
        MSO_SHAPE_TYPE.CHART: "Chart",
        MSO_SHAPE_TYPE.COMMENT: "Comment",
        MSO_SHAPE_TYPE.CONNECTOR: "Connector",
        MSO_SHAPE_TYPE.DIAGRAM: "Diagram",
        MSO_SHAPE_TYPE.FREE_FORM: "FreeForm",
        MSO_SHAPE_TYPE.GROUP: "Group",
        MSO_SHAPE_TYPE.INK: "Ink",
        MSO_SHAPE_TYPE.LINE: "Line",
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: "LinkedOLE",
        MSO_SHAPE_TYPE.LINKED_PICTURE: "LinkedPicture",
        MSO_SHAPE_TYPE.MEDIA: "Media",
        MSO_SHAPE_TYPE.OLE_OBJECT: "OLEObject",
        MSO_SHAPE_TYPE.PICTURE: "Picture",
        MSO_SHAPE_TYPE.PLACEHOLDER: "Placeholder",
        MSO_SHAPE_TYPE.RADIO_BUTTON: "RadioButton",
        MSO_SHAPE_TYPE.SCRIPT_ANCHOR: "ScriptAnchor",
        MSO_SHAPE_TYPE.SHAPE_TYPE_MIX: "ShapeTypeMix",
        MSO_SHAPE_TYPE.TABLE: "Table",
        MSO_SHAPE_TYPE.TEXT_BOX: "TextBox",
        MSO_SHAPE_TYPE.TEXT_EFFECT: "TextEffect",
        MSO_SHAPE_TYPE.WEB_VIDEO: "WebVideo",
    }
    return mapping.get(st, f"Unknown({st})")

def is_chinese_font(font_name):
    """Heuristic: check if font name suggests Chinese support."""
    if not font_name:
        return False
    name_lower = font_name.lower()
    chinese_keywords = [
        "hei", "song", "fang", "kai", "ming", "yuan", "wenquan", "noto sans sc",
        "noto serif sc", "pingfang", "microsoft yahei", "simsun", "simhei",
        "fz", "fangzheng", "华文", "微软雅黑", "黑体", "宋体", "楷体", "仿宋",
        "思源", "汉仪", "方正"
    ]
    for kw in chinese_keywords:
        if kw in name_lower:
            return True
    # Also check if it's a common CJK font
    # SimSun, SimHei, etc.
    if any(c in name_lower for c in ['sim', 'cjk', 'chinese', 'sc ', ' sc', 'cn']):
        return True
    return False

def main():
    results = {
        "file_path": PPT_PATH,
        "zip_valid": True,
        "total_slides": 0,
        "slides": [],
        "notes": {},
        "fonts": {"all": [], "chinese": set(), "non_chinese": set()},
        "shape_stats": Counter(),
        "layout_stats": Counter(),
        "errors": [],
    }

    # 1. ZIP pre-check
    import zipfile
    results["zip_valid"] = zipfile.is_zipfile(PPT_PATH)

    if not results["zip_valid"]:
        results["errors"].append("FATAL: File is not a valid ZIP (not a valid PPTX)")
        print(json.dumps(results, ensure_ascii=False, indent=2, default=str))
        return

    # 2. Open presentation
    try:
        prs = Presentation(PPT_PATH)
    except Exception as e:
        results["errors"].append(f"FATAL: Cannot open Presentation: {e}")
        print(json.dumps(results, ensure_ascii=False, indent=2, default=str))
        return

    results["total_slides"] = len(prs.slides)

    # 3. Process each slide
    for idx, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": idx,
            "layout": None,
            "text": "",
            "text_preview": "",
            "shape_types": [],
            "shapes_detail": [],
            "has_notes": False,
            "notes_text": "",
        }
        try:
            # Layout
            try:
                slide_data["layout"] = slide.slide_layout.name if slide.slide_layout else "Unknown"
            except:
                slide_data["layout"] = "Error getting layout"

            # Shapes
            for shape in slide.shapes:
                st = shape_type_name(shape.shape_type)
                slide_data["shape_types"].append(st)

                shape_info = {"name": shape.name, "type": st, "has_text": False, "text_preview": ""}
                
                # Extract text
                if shape.has_text_frame:
                    try:
                        text = shape.text.strip()
                        if text:
                            shape_info["has_text"] = True
                            shape_info["text_preview"] = text[:80]
                            slide_data["text"] += text + "\n"
                    except:
                        pass

                # Font audit on text runs
                if shape.has_text_frame:
                    try:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    fn = run.font.name
                                    if fn:
                                        results["fonts"]["all"].append(fn)
                                        if is_chinese_font(fn):
                                            results["fonts"]["chinese"].add(fn)
                                        else:
                                            results["fonts"]["non_chinese"].add(fn)
                                except:
                                    pass
                    except:
                        pass

                slide_data["shapes_detail"].append(shape_info)

            # Update shape stats (global)
            for st in slide_data["shape_types"]:
                results["shape_stats"][st] += 1

            # Update layout stats
            results["layout_stats"][slide_data["layout"]] += 1

            # Text preview
            slide_data["text_preview"] = slide_data["text"][:100].replace("\n", " | ")

            # Speaker notes
            try:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide and notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_data["has_notes"] = True
                            slide_data["notes_text"] = notes_text
                            results["notes"][idx] = notes_text
            except:
                pass

        except Exception as e:
            results["errors"].append(f"Slide {idx}: {e}")

        results["slides"].append(slide_data)

    # Convert sets to lists for JSON
    results["fonts"]["chinese"] = sorted(results["fonts"]["chinese"])
    results["fonts"]["non_chinese"] = sorted(results["fonts"]["non_chinese"])
    results["font_total_unique"] = len(set(results["fonts"]["all"]))

    # Print as JSON
    print(json.dumps(results, ensure_ascii=False, indent=2, default=str))

if __name__ == "__main__":
    main()
